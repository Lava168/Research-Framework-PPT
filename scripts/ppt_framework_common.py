"""Shared layout primitives for Research Framework PPT (structure only).

Geometry targets 16:9 widescreen (13.333 x 7.5 in). Callers supply labels;
this module draws boxes, arrows, and grids — no domain text.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Canvas
# ---------------------------------------------------------------------------
SLIDE_W = 13.333
SLIDE_H = 7.5

# ---------------------------------------------------------------------------
# Palette (matches Research Framework.pptx)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x12, 0x2A, 0x52)
NAVY2 = RGBColor(0x0F, 0x2B, 0x55)
BLUE = RGBColor(0x1B, 0x5A, 0xA6)
BORDER = RGBColor(0x7C, 0x93, 0xB8)
BORDER2 = RGBColor(0x7C, 0x9C, 0xC8)
PANEL_BORDER = RGBColor(0xB9, 0xBE, 0xC8)
GRAY_BORDER = RGBColor(0x8A, 0x8F, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE9, 0xEF, 0xF7)
LIGHT2 = RGBColor(0xDC, 0xE6, 0xF2)
LIGHT3 = RGBColor(0xC7, 0xD6, 0xEC)
GREEN_HDR = RGBColor(0xDD, 0xEA, 0xD6)
YELLOW = RGBColor(0xFB, 0xF2, 0xD3)
DARK = RGBColor(0x20, 0x2A, 0x3A)
GRAY = RGBColor(0x55, 0x5F, 0x70)
ARROW = RGBColor(0x3A, 0x4C, 0x6B)
ACCENT_BAR = (
    RGBColor(0x1B, 0x5A, 0xA6),
    RGBColor(0x6D, 0x4B, 0xA3),
    RGBColor(0xB8, 0x5C, 0x00),
    RGBColor(0x0F, 0x6B, 0x3B),
    RGBColor(0x0E, 0x6E, 0x86),
)
MID_BLUE = RGBColor(0x2F, 0x64, 0xA7)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_dash(shape, val: str = "dash") -> None:
    ln = shape.line._get_or_add_ln()
    for e in ln.findall(qn("a:prstDash")):
        ln.remove(e)
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": val}))


def rect(slide, x, y, w, h, fill, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    return sp


def label(slide, x, y, w, h, text, *, size=9, bold=False, color=DARK,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Arial"):
    if not text:
        return None
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    return box


def header_bar(slide, x, y, w, h, text="", *, fill=LIGHT, color=NAVY, size=10,
               border=NAVY, line_w=1.0):
    rect(slide, x, y, w, h, fill, line=border, line_w=line_w)
    if text:
        label(slide, x, y, w, h, text, size=size, bold=True, color=color)


def panel_outline(slide, x, y, w, h, *, border=PANEL_BORDER):
    rect(slide, x, y, w, h, None, line=border, line_w=1.25)


def card(slide, x, y, w, h, *, fill=WHITE, border=BORDER, line_w=1.0,
         dash=False, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = rect(slide, x, y, w, h, fill, line=border, line_w=line_w, shape=shape)
    if dash:
        set_dash(sp)
    return sp


def arrow(slide, x, y, w, h, *, shape=MSO_SHAPE.RIGHT_ARROW, color=BLUE):
    return rect(slide, x, y, w, h, color, shape=shape)


def stacked_group(slide, x, y, w, h, header, n_sub, *, header_h=0.26, pad=0.07,
                  sub_gap=0.08, footer_h=0.0, header_fill=LIGHT, border=NAVY):
    """Outer box + header strip + N equal sub-cards (+ optional footer band)."""
    rect(slide, x, y, w, h, WHITE, line=border, line_w=1.25)
    header_bar(slide, x, y, w, header_h, header, fill=header_fill, border=border)
    body_top = y + header_h + pad
    body_h = h - header_h - 2 * pad - footer_h
    if n_sub <= 0:
        return
    sw = (w - 2 * pad - sub_gap * (n_sub - 1)) / n_sub
    for i in range(n_sub):
        card(slide, x + pad + i * (sw + sub_gap), body_top, sw, body_h, border=BORDER)
    if footer_h > 0:
        card(slide, x + pad, y + h - footer_h - pad * 0.3, w - 2 * pad, footer_h,
             border=BORDER, fill=WHITE)


def flow_row(slide, x, y, w, h, n_steps, *, gap=0.33, arrow_w=0.29, arrow_h=0.18,
             box_border=BORDER):
    """Horizontal flow: box → arrow → box → … (n_steps boxes)."""
    if n_steps < 1:
        return
    total_arrow = arrow_w * max(0, n_steps - 1)
    total_gap = gap * max(0, n_steps - 1)
    bw = (w - total_arrow - total_gap) / n_steps
    cx = x
    for i in range(n_steps):
        card(slide, cx, y, bw, h, border=box_border)
        cx += bw
        if i < n_steps - 1:
            ay = y + h / 2 - arrow_h / 2
            arrow(slide, cx + (gap - arrow_w) / 2, ay, arrow_w, arrow_h)
            cx += gap


def chapter_strip(slide, x, y, w, h, chapter_tag="", name="", *, tag_w=1.0):
    """Rounded container with tag + name row (structure only)."""
    card(slide, x, y, w, h, rounded=True, border=BORDER)
    if chapter_tag:
        rect(slide, x + 0.17, y + 0.10, tag_w, 0.26, NAVY)
        label(slide, x + 0.17, y + 0.10, tag_w, 0.26, chapter_tag,
              size=8.5, bold=True, color=WHITE)
    if name:
        label(slide, x + 0.17 + tag_w + 0.12, y + 0.10, 3.0, 0.26, name,
              size=11, bold=True, color=NAVY, align=PP_ALIGN.LEFT)


def numbered_steps(slide, x_items, y, w_item, h_item, n, *, start=1):
    """Vertical numbered circles with empty text bands."""
    for i in range(n):
        oy = y + i * (h_item + 0.08)
        cx, cy = x_items, oy + 0.06
        rect(slide, cx, cy, 0.27, 0.27, WHITE, line=BLUE, line_w=1.2,
             shape=MSO_SHAPE.OVAL)
        label(slide, cx, cy, 0.27, 0.27, str(start + i), size=8, bold=True, color=BLUE)
        card(slide, cx + 0.38, oy, w_item - 0.38, h_item, border=BORDER2)


def accent_cards_row(slide, x, y, w, h, n, *, gap=0.11):
    """N equal cards with colored top accent bars."""
    cw = (w - gap * (n - 1)) / n
    for i in range(n):
        cx = x + i * (cw + gap)
        card(slide, cx, y, cw, h, rounded=True, border=BORDER2)
        rect(slide, cx, y, cw, 0.073, ACCENT_BAR[i % len(ACCENT_BAR)])


def save_prs(prs: Presentation, path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path
